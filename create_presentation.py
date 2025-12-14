#!/usr/bin/env python3
"""
Script to create a PowerPoint presentation for UniChat project.
Requires: pip install python-pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    # Create presentation object
    prs = Presentation()
    
    # Set slide dimensions (16:9)
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1 - UniChat Overview
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Title
    title_box = slide1.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "UniChat Overview"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0, 51, 102)
    
    # Content
    content_box = slide1.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    content_text = """• One codebase, three platforms (iOS/Android/Web) with Ionic + Capacitor.
• Simple multilingual chat that feels native everywhere.
• Visuals: onboarding and language choices (assets/onboarding.png, assets/translation-preferences.png)."""
    
    content_frame.text = content_text
    for para in content_frame.paragraphs:
        para.font.size = Pt(18)
        para.space_after = Pt(12)
        para.level = 0
    
    # Slide 2 - Problem & Rationale
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    
    title_box2 = slide2.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame2 = title_box2.text_frame
    title_frame2.text = "Problem & Rationale"
    title_para2 = title_frame2.paragraphs[0]
    title_para2.font.size = Pt(44)
    title_para2.font.bold = True
    title_para2.font.color.rgb = RGBColor(0, 51, 102)
    
    content_box2 = slide2.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
    content_frame2 = content_box2.text_frame
    content_frame2.word_wrap = True
    
    content_text2 = """• Chat translation is often an afterthought, slow, or awkward.
• Two native apps mean double work and drifting designs.
• Use a hybrid Ionic app with a small client-side translator to stay fast and consistent."""
    
    content_frame2.text = content_text2
    for para in content_frame2.paragraphs:
        para.font.size = Pt(20)
        para.space_after = Pt(14)
    
    # Slide 3 - Architecture
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    
    title_box3 = slide3.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame3 = title_box3.text_frame
    title_frame3.text = "Architecture"
    title_para3 = title_frame3.paragraphs[0]
    title_para3.font.size = Pt(44)
    title_para3.font.bold = True
    title_para3.font.color.rgb = RGBColor(0, 51, 102)
    
    content_box3 = slide3.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
    content_frame3 = content_box3.text_frame
    content_frame3.word_wrap = True
    
    content_text3 = """• Ionic React UI (IonPage/IonHeader/IonContent) with platform-aware styling.
• Translation broker on the client; primary provider with a simple fallback and cache.
• Messages keep both original and translated text; offline queue retries when back online.
• Capacitor plugins for secure storage, clipboard, and network status; native shells via ionic cap build.
• Visual: login (assets/fig1.png) or conversation list (assets/fig2.png)."""
    
    content_frame3.text = content_text3
    for para in content_frame3.paragraphs:
        para.font.size = Pt(18)
        para.space_after = Pt(12)
    
    # Slide 4 - Method in Brief
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    
    title_box4 = slide4.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame4 = title_box4.text_frame
    title_frame4.text = "Method in Brief"
    title_para4 = title_frame4.paragraphs[0]
    title_para4.font.size = Pt(44)
    title_para4.font.bold = True
    title_para4.font.color.rgb = RGBColor(0, 51, 102)
    
    content_box4 = slide4.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
    content_frame4 = content_box4.text_frame
    content_frame4.word_wrap = True
    
    content_text4 = """• Small bilingual group across four language pairs.
• Tasks: onboard, set languages, start a chat, send and read translated messages, go offline then back online.
• Logs: translation and delivery speed, cache use, stability, network quality; short usability surveys.
• Devices: iPhone, Android phone, and web browser from the same codebase."""
    
    content_frame4.text = content_text4
    for para in content_frame4.paragraphs:
        para.font.size = Pt(18)
        para.space_after = Pt(12)
    
    # Slide 5 - Results
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    
    title_box5 = slide5.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame5 = title_box5.text_frame
    title_frame5.text = "Results"
    title_para5 = title_frame5.paragraphs[0]
    title_para5.font.size = Pt(44)
    title_para5.font.bold = True
    title_para5.font.color.rgb = RGBColor(0, 51, 102)
    
    content_box5 = slide5.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
    content_frame5 = content_box5.text_frame
    content_frame5.word_wrap = True
    
    content_text5 = """• Translation stayed quick and accurate enough for casual chat.
• Experience felt the same on iOS, Android, and web.
• Worked even on slow networks thanks to caching and retry.
• Users liked the language badges and offline queue; no major crashes noted.
• Visuals: chat detail (assets/fig3.png), analytics cards (assets/fig4.png), runtime dashboard (assets/performance-dashboard.png)."""
    
    content_frame5.text = content_text5
    for para in content_frame5.paragraphs:
        para.font.size = Pt(16)
        para.space_after = Pt(10)
    
    # Slide 6 - Conclusion
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    
    title_box6 = slide6.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame6 = title_box6.text_frame
    title_frame6.text = "Conclusion"
    title_para6 = title_frame6.paragraphs[0]
    title_para6.font.size = Pt(44)
    title_para6.font.bold = True
    title_para6.font.color.rgb = RGBColor(0, 51, 102)
    
    content_box6 = slide6.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
    content_frame6 = content_box6.text_frame
    content_frame6.word_wrap = True
    
    content_text6 = """• Hybrid Ionic + client translation met the goals without building two native apps.
• Next: try on-device translation, smarter caching for bad networks, better RTL/tone support, and add end-to-end encryption."""
    
    content_frame6.text = content_text6
    for para in content_frame6.paragraphs:
        para.font.size = Pt(20)
        para.space_after = Pt(14)
    
    # Save presentation
    output_file = "UniChat_Presentation.pptx"
    prs.save(output_file)
    print(f"✓ Presentation created successfully: {output_file}")
    print(f"  Location: {output_file}")
    return output_file

if __name__ == "__main__":
    try:
        create_presentation()
    except ImportError:
        print("Error: python-pptx library not found.")
        print("Please install it using: pip install python-pptx")
        exit(1)
    except Exception as e:
        print(f"Error creating presentation: {e}")
        exit(1)

